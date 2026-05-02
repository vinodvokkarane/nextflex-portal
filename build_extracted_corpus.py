"""
One-time pre-extraction script.

Run this LOCALLY (not on Render) to produce public_data/extracted_corpus.json
which contains all the PDF/PPTX text already extracted. The deployed server
then just reads that JSON file — no pypdf, no python-pptx at runtime, no
memory spikes.

Usage:
    python3 build_extracted_corpus.py

Output:
    public_data/extracted_corpus.json   (~5-10 MB JSON file with all text)

Run this ONCE whenever public_data/ files change. Commit the JSON to git
so Render can use it directly.
"""

import csv
import json
import re
import sys
from pathlib import Path

ROOT = Path(__file__).parent
DATA_ROOT = ROOT / "public_data"
PAPERS_DIR = DATA_ROOT / "nextflex_acknowledged_papers_2016_2026"
OUTPUT = DATA_ROOT / "extracted_corpus.json"


def extract_pdf(path: Path, max_pages: int = 40, max_chars: int = 200_000):
    """Stream-extract PDF text page by page."""
    import pypdf
    text_parts = []
    total = 0
    pages_total = 0
    try:
        with open(path, "rb") as fh:
            reader = pypdf.PdfReader(fh)
            pages_total = len(reader.pages)
            for i in range(min(pages_total, max_pages)):
                if total >= max_chars:
                    break
                try:
                    t = reader.pages[i].extract_text() or ""
                except Exception:
                    t = ""
                if t:
                    text_parts.append(t)
                    total += len(t)
    except Exception as e:
        print(f"  ! PDF extract failed for {path.name}: {e}", file=sys.stderr)
        return "", 0
    return "\n".join(text_parts), pages_total


def extract_pptx(path: Path):
    from pptx import Presentation
    try:
        prs = Presentation(str(path))
        chunks = []
        n_slides = 0
        for slide in prs.slides:
            n_slides += 1
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            slide_text.append(t)
            if slide_text:
                chunks.append(f"[Slide {n_slides}]\n" + "\n".join(slide_text))
        return "\n\n".join(chunks), n_slides
    except Exception as e:
        print(f"  ! PPTX extract failed for {path.name}: {e}", file=sys.stderr)
        return "", 0


def main():
    if not DATA_ROOT.exists():
        print(f"public_data/ not found at {DATA_ROOT}", file=sys.stderr)
        sys.exit(1)

    corpus = {}

    # Papers
    papers_manifest = PAPERS_DIR / "manifest.csv"
    if papers_manifest.exists():
        with papers_manifest.open(newline="", encoding="utf-8") as f:
            for row in csv.DictReader(f):
                folder = (row.get("folder") or "").strip()
                filename = (row.get("filename") or "").strip()
                if not filename:
                    continue
                rel = f"nextflex_acknowledged_papers_2016_2026/{folder}/{filename}"
                local = PAPERS_DIR / folder / filename
                if not local.exists():
                    continue
                print(f"  paper: {filename[:60]}")
                text, pages = extract_pdf(local)
                corpus[rel] = {"text": text, "pages": pages}

    # Funded assets
    funded_csv = DATA_ROOT / "nextflex_assets_included_public_files.csv"
    if funded_csv.exists():
        with funded_csv.open(newline="", encoding="utf-8") as f:
            for row in csv.DictReader(f):
                rel = (row.get("downloaded_filename") or "").strip()
                if not rel:
                    continue
                local = DATA_ROOT / rel
                if not local.exists():
                    continue
                print(f"  {row.get('category', '?'):8} {Path(rel).name[:60]}")
                if rel.lower().endswith(".pdf"):
                    text, pages = extract_pdf(local)
                elif rel.lower().endswith(".pptx"):
                    text, pages = extract_pptx(local)
                else:
                    continue
                corpus[rel] = {"text": text, "pages": pages}

    OUTPUT.write_text(json.dumps(corpus, ensure_ascii=False, indent=None))
    size_mb = OUTPUT.stat().st_size / 1024 / 1024
    total_chars = sum(len(v["text"]) for v in corpus.values())
    print()
    print(f"Wrote {OUTPUT} ({size_mb:.2f} MB)")
    print(f"  {len(corpus)} files, {total_chars:,} chars total")


if __name__ == "__main__":
    main()
