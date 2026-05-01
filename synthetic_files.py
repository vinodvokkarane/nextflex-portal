"""
Synthetic file generator for NextFlex projects.

For each project in the database, generates:
  - A PDF "Final Report" (8-12 pages) with sections grounded in the ontology
  - A PPTX "Project Briefing" (8 slides) with the same content structured visually

Both use the EXACT taxonomy from the OWL-DL ontology shown in the manager
dashboard:
  Materials: Conductive inks, Dielectric inks, Substrates, Active components
  Processes: Sintering, UV cure, Print parameters, Assembly
  Performance: S-parameters, Power & efficiency, Physical & EM

Generated content gets indexed into the chunks table so GraphRAG retrieval
hits the actual file contents.

Files are written to /tmp/nfx-synthetic-files/ (ephemeral on Render free
tier; rebuild on cold start in ~30 seconds for 390 projects).
"""

import json
import os
import sqlite3
import tempfile
from pathlib import Path
from io import BytesIO

from reportlab.lib import colors
from reportlab.lib.pagesizes import LETTER
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak,
)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Ontology (matches OWL-DL block in manager.html exactly)
ONTOLOGY = {
    "Materials": {
        "Conductive inks": ["vendor", "formulation", "solids fraction",
                            "dilution solvent", "post-cure resistivity"],
        "Dielectric inks": ["vendor", "cure mechanism", "loss tangent",
                            "dielectric constant", "Rq"],
        "Substrates": ["material", "thickness", "CTE", "Dk"],
        "Active components": ["part number", "technology node",
                              "frequency range", "power class"],
    },
    "Processes": {
        "Sintering": ["protocol stages", "ramp rates",
                      "dwell temperatures", "dwell times"],
        "UV cure": ["intensity", "wavelength", "duration", "in-situ / ex-situ"],
        "Print parameters": ["nozzle size", "speed", "flowrates",
                             "platen temperature"],
        "Assembly": ["die attach method", "epoxy type", "cure schedule"],
    },
    "Performance metrics": {
        "S-parameters @ frequency": ["gain", "return loss",
                                     "insertion loss", "isolation"],
        "Power & efficiency": ["output power", "PAE",
                               "CW power handling", "pulsed power handling"],
        "Physical & EM": ["surface roughness Rq", "dimensional tolerances",
                          "EMI isolation", "bandwidth"],
    },
}

OUTPUT_DIR = Path(os.environ.get("NFX_FILES_DIR", "/tmp/nfx-synthetic-files"))
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


def _get_db_path():
    return Path(__file__).parent / "nextflex.db"


def _classify_material(material_name: str) -> str:
    """Map a material name to its ontology subclass."""
    n = material_name.lower()
    if any(k in n for k in ("kapton", "pen substrate", "tpu", "pet", "glass-pi")):
        return "Substrates"
    if any(k in n for k in ("dielectric", "bst", "el-p", "5018", "cr-18")):
        return "Dielectric inks"
    if any(k in n for k in ("mmic", "ad8009", "lmv", "qorvo", "tga")):
        return "Active components"
    return "Conductive inks"


def _classify_process(process_name: str) -> str:
    n = process_name.lower()
    if "sinter" in n:
        return "Sintering"
    if "uv cure" in n:
        return "UV cure"
    if "die attach" in n or "wire bond" in n or "anisotropic" in n:
        return "Assembly"
    if "etch" in n:
        return "Print parameters"
    return "Print parameters"


def _classify_performance(perf_name: str) -> str:
    n = perf_name.lower()
    if "power amplifier" in n or "pae" in n:
        return "Power & efficiency"
    if "roughness" in n or "sheet resistance" in n or "thermal cycle" in n:
        return "Physical & EM"
    return "S-parameters @ frequency"


def _props_table(props: dict, headers=("Property", "Value")):
    """Build a property table for ReportLab."""
    if not props:
        return None
    rows = [headers] + [(k.replace('_', ' '), str(v)) for k, v in props.items()]
    t = Table(rows, colWidths=[2.6 * inch, 3.2 * inch])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#0A8F8F')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 9),
        ('PADDING', (0, 0), (-1, -1), 7),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#cbd5e1')),
        ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#f8fafc')),
        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
    ]))
    return t


def generate_pdf(project: dict, materials: list, processes: list,
                 perfs: list) -> bytes:
    """Generate a multi-page PDF report for a project."""
    buf = BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=LETTER,
        rightMargin=0.7*inch, leftMargin=0.7*inch,
        topMargin=0.6*inch, bottomMargin=0.6*inch,
        title=project["title"], author=project["lead_institution"],
    )

    styles = getSampleStyleSheet()
    h1 = ParagraphStyle('H1', parent=styles['Heading1'],
                        textColor=colors.HexColor('#132B4A'),
                        fontSize=18, spaceAfter=10, leading=22)
    h2 = ParagraphStyle('H2', parent=styles['Heading2'],
                        textColor=colors.HexColor('#0A8F8F'),
                        fontSize=13, spaceAfter=6, spaceBefore=14)
    body = ParagraphStyle('Body', parent=styles['Normal'],
                          fontSize=10.5, leading=15, spaceAfter=8)
    meta = ParagraphStyle('Meta', parent=styles['Normal'],
                          fontSize=9, textColor=colors.HexColor('#64748b'),
                          spaceAfter=14)
    cap = ParagraphStyle('Cap', parent=styles['Normal'],
                         fontSize=8, textColor=colors.HexColor('#64748b'),
                         spaceAfter=14, alignment=1)

    story = []

    # Cover
    story.append(Paragraph(
        f"<b>NextFlex Final Report &mdash; {project['project_call']}</b>", meta))
    story.append(Paragraph(project["title"], h1))
    pis = ", ".join(project.get("principal_investigators", [])) or "TBD"
    partners = ", ".join(project.get("industry_partners", [])) or "None"
    story.append(Paragraph(
        f"<b>Lead institution:</b> {project['lead_institution']}<br/>"
        f"<b>Principal investigators:</b> {pis}<br/>"
        f"<b>Industry partners:</b> {partners}<br/>"
        f"<b>Period of performance:</b> {project['start_date']} to {project['end_date']}<br/>"
        f"<b>Funding:</b> ${project['funding_amount']:,}<br/>"
        f"<b>Focus area:</b> {project['focus_area']}<br/>"
        f"<b>Congressional district:</b> {project['congressional_district']}<br/>"
        f"<b>TRL:</b> {project['trl_start']} &rarr; {project['trl_end']}<br/>"
        f"<b>Classification:</b> {project['classification']}", body))

    # 1. Executive Summary
    story.append(Paragraph("1. Executive Summary", h2))
    story.append(Paragraph(project["abstract"], body))

    # 2. Materials section — ontology subclass-driven
    story.append(Paragraph("2. Materials", h2))
    story.append(Paragraph(
        "Materials selection followed the NextFlex OWL-DL ontology, with each candidate validated "
        "against vendor specifications and prior NextFlex Project Call data. The ontology defines "
        "four sub-classes of materials: Conductive inks, Dielectric inks, Substrates, "
        "and Active components. Selections for this project are documented below.", body))

    for mat in materials[:4]:
        subclass = _classify_material(mat["name"])
        story.append(Paragraph(f"<b>2.x {mat['name']}</b> &mdash; "
                                f"<i>ontology subclass: {subclass}</i>", body))
        if mat.get("properties"):
            t = _props_table(mat["properties"])
            if t:
                story.append(t)
                story.append(Paragraph(
                    f"Table: properties of {mat['name']} per vendor datasheet "
                    f"and PERC/RURI characterization.", cap))

    story.append(PageBreak())

    # 3. Processes
    story.append(Paragraph("3. Process & Manufacturing", h2))
    story.append(Paragraph(
        "Process parameters were optimized through a design-of-experiments approach. "
        "The OWL-DL ontology defines four process sub-classes: Sintering, UV cure, "
        "Print parameters, and Assembly. Final operating windows are documented below.", body))

    for proc in processes[:3]:
        subclass = _classify_process(proc["name"])
        story.append(Paragraph(f"<b>3.x {proc['name']}</b> &mdash; "
                                f"<i>ontology subclass: {subclass}</i>", body))
        if proc.get("properties"):
            t = _props_table(proc["properties"])
            if t:
                story.append(t)
                story.append(Paragraph(
                    f"Table: optimized {subclass.lower()} parameters for {proc['name']}.",
                    cap))

    story.append(PageBreak())

    # 4. Results
    story.append(Paragraph("4. Performance Results", h2))
    story.append(Paragraph(
        "Performance characterization used the three OWL-DL performance sub-classes: "
        "S-parameters @ frequency, Power & efficiency, and Physical & EM. "
        "Statistical analysis across multiple independent samples confirms repeatability "
        "within design tolerances.", body))

    for perf in perfs[:3]:
        subclass = _classify_performance(perf["name"])
        story.append(Paragraph(f"<b>4.x {perf['name']}</b> &mdash; "
                                f"<i>ontology subclass: {subclass}</i>", body))
        if perf.get("properties"):
            t = _props_table(perf["properties"], headers=("Metric", "Measured"))
            if t:
                story.append(t)
                story.append(Paragraph(
                    f"Table: measured {subclass.lower()} for the integrated demonstrator.", cap))

    # 5. Reliability
    story.append(Paragraph("5. Reliability & Test", h2))
    story.append(Paragraph(
        "Reliability characterization included thermal cycling per IPC-9701 "
        "(-40\u00B0C to +125\u00B0C, 500 cycles), humidity exposure per JEDEC "
        "JESD22-A101 (85\u00B0C/85% RH, 1000 hours), and mechanical flex testing "
        "(50,000 cycles at 10 mm bend radius). All samples met the program "
        "qualification criteria with margin.", body))

    # 6. Conclusions
    story.append(Paragraph("6. Conclusions & Transition Pathways", h2))
    story.append(Paragraph(project["outcomes"], body))
    story.append(Paragraph(
        f"Results from this {project['project_call']} effort have been entered into "
        f"the NextFlex Secure AI/ML Knowledge Base and are queryable via the GraphRAG "
        f"interface. Industry partners {partners} are identified collaborators for "
        f"transition to TRL-7+.", body))

    # Citations
    pubs = project.get("publications", [])
    pats = project.get("patents", [])
    if pubs or pats:
        story.append(Paragraph("7. Publications & Patents", h2))
        for pub in pubs:
            story.append(Paragraph(f"&bull; {pub}", body))
        for pat in pats:
            story.append(Paragraph(f"&bull; {pat}", body))

    doc.build(story)
    return buf.getvalue()


def generate_pptx(project: dict, materials: list, processes: list,
                  perfs: list) -> bytes:
    """Generate a PPTX project briefing."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    NAVY = RGBColor(0x13, 0x2B, 0x4A)
    TEAL = RGBColor(0x0A, 0x8F, 0x8F)
    SLATE = RGBColor(0x1f, 0x29, 0x37)
    MUTED = RGBColor(0x64, 0x74, 0x8b)
    LIGHT = RGBColor(0xf8, 0xfa, 0xfc)

    def add_title_bar(slide, text, sub=None):
        # Title strip
        from pptx.util import Emu
        from pptx.shapes.connector import Connector
        # Title text
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4),
                                       Inches(12), Inches(0.7))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = NAVY
        if sub:
            tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(1.1),
                                           Inches(12), Inches(0.4))
            p2 = tb2.text_frame.paragraphs[0]
            p2.text = sub
            p2.font.size = Pt(13); p2.font.color.rgb = TEAL; p2.font.bold = True

    def add_bullets(slide, top_inches, items, bullet_size=14, leading=1.4):
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(top_inches),
                                       Inches(12), Inches(5))
        tf = tb.text_frame; tf.word_wrap = True
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "\u2022  " + it
            p.font.size = Pt(bullet_size); p.font.color.rgb = SLATE
            p.space_after = Pt(8)

    def add_kv_block(slide, top_inches, kv: dict, title: str = None):
        if title:
            tb = slide.shapes.add_textbox(Inches(0.6), Inches(top_inches),
                                           Inches(12), Inches(0.4))
            p = tb.text_frame.paragraphs[0]
            p.text = title
            p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = TEAL
            top_inches += 0.45
        for k, v in kv.items():
            tb = slide.shapes.add_textbox(Inches(0.6), Inches(top_inches),
                                           Inches(12), Inches(0.32))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.text = f"{k}: "
            p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = MUTED
            r = p.add_run(); r.text = str(v)
            r.font.size = Pt(11); r.font.bold = False; r.font.color.rgb = SLATE
            top_inches += 0.34

    blank = prs.slide_layouts[6]

    # Slide 1: Cover
    s = prs.slides.add_slide(blank)
    tb = s.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = f"NEXTFLEX  \u00B7  {project['project_call']}  \u00B7  PROJECT BRIEFING"
    p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = TEAL

    tb = s.shapes.add_textbox(Inches(0.6), Inches(2.3), Inches(12), Inches(2))
    p = tb.text_frame.paragraphs[0]
    p.text = project["title"]
    p.font.size = Pt(32); p.font.bold = True; p.font.color.rgb = NAVY

    tb = s.shapes.add_textbox(Inches(0.6), Inches(4.4), Inches(12), Inches(2))
    tf = tb.text_frame
    info = [
        f"Lead: {project['lead_institution']}",
        f"PIs: {', '.join(project.get('principal_investigators', []))}",
        f"Period: {project['start_date']} to {project['end_date']}",
        f"Funding: ${project['funding_amount']:,}",
        f"Focus: {project['focus_area']}  \u00B7  District: {project['congressional_district']}",
        f"TRL: {project['trl_start']} \u2192 {project['trl_end']}",
    ]
    for i, line in enumerate(info):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(14); p.font.color.rgb = SLATE
        p.space_after = Pt(4)

    # Slide 2: Executive summary
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Executive Summary", "Project objective and key outcomes")
    add_bullets(s, 1.9, [
        project["abstract"][:300],
        project["outcomes"],
    ], bullet_size=14)

    # Slide 3: Materials (ontology-classified)
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Materials", "Ontology-classified per OWL-DL schema")
    bullets = []
    for mat in materials[:4]:
        subclass = _classify_material(mat["name"])
        bullets.append(f"{mat['name']}  \u2014  [{subclass}]")
        if mat.get("properties"):
            kv = ", ".join(f"{k}={v}" for k, v in
                           list(mat["properties"].items())[:3])
            bullets.append(f"     {kv}")
    add_bullets(s, 1.9, bullets, bullet_size=12)

    # Slide 4: Processes
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Process & Manufacturing",
                  "Print + post-process + assembly")
    bullets = []
    for proc in processes[:3]:
        subclass = _classify_process(proc["name"])
        bullets.append(f"{proc['name']}  \u2014  [{subclass}]")
        if proc.get("properties"):
            kv = ", ".join(f"{k}={v}" for k, v in
                           list(proc["properties"].items())[:4])
            bullets.append(f"     {kv}")
    add_bullets(s, 1.9, bullets, bullet_size=12)

    # Slide 5: Performance
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Performance Results",
                  "S-parameters, Power, Physical & EM")
    bullets = []
    for perf in perfs[:3]:
        subclass = _classify_performance(perf["name"])
        bullets.append(f"{perf['name']}  \u2014  [{subclass}]")
        if perf.get("properties"):
            kv = ", ".join(f"{k}={v}" for k, v in
                           list(perf["properties"].items())[:4])
            bullets.append(f"     {kv}")
    add_bullets(s, 1.9, bullets, bullet_size=12)

    # Slide 6: Reliability
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Reliability & Test", "Environmental qualification")
    add_bullets(s, 1.9, [
        "Thermal cycling per IPC-9701: -40\u00B0C to +125\u00B0C, 500 cycles",
        "Humidity per JEDEC JESD22-A101: 85\u00B0C / 85% RH, 1000 hours",
        "Mechanical flex: 50,000 cycles @ 10 mm bend radius",
        "All samples met program qualification criteria with margin.",
    ], bullet_size=14)

    # Slide 7: Transition pathway
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Transition Pathway", "Commercialization & DoD acquisition")
    pathway = [
        f"Industry partners: {', '.join(project.get('industry_partners', [])) or 'TBD'}",
        f"PEO alignment: {project.get('peo') or 'TBD'}",
        f"Classification: {project['classification']}",
        "Results entered into NextFlex Secure AI/ML Knowledge Base",
        "Queryable via GraphRAG interface for downstream programs",
    ]
    add_bullets(s, 1.9, pathway, bullet_size=14)

    # Slide 8: References
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Publications & Patents", "Peer-reviewed outputs")
    refs = [pub for pub in project.get("publications", [])]
    refs += [pat for pat in project.get("patents", [])]
    if not refs:
        refs = ["No publications or patents filed during the period of performance."]
    add_bullets(s, 1.9, refs[:8], bullet_size=11)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def get_or_generate_files(project_id: str):
    """Generate (and cache) the synthetic PDF + PPTX for a project. Returns paths."""
    pdf_path = OUTPUT_DIR / f"{project_id}.pdf"
    pptx_path = OUTPUT_DIR / f"{project_id}.pptx"

    if pdf_path.exists() and pptx_path.exists():
        return pdf_path, pptx_path

    # Pull project + entity data from DB
    db = sqlite3.connect(_get_db_path())
    db.row_factory = sqlite3.Row
    prow = db.execute("SELECT * FROM projects WHERE id = ?", (project_id,)).fetchone()
    if not prow:
        db.close()
        raise ValueError(f"Project {project_id} not found")

    project = dict(prow)
    for f in ("principal_investigators", "co_investigators", "industry_partners",
              "materials_used", "processes_used", "publications", "patents", "keywords"):
        try:
            project[f] = json.loads(project.get(f) or "[]")
        except Exception:
            project[f] = []

    # Get the actual entity records for materials/processes/perfs used in this project
    rels = db.execute(
        "SELECT DISTINCT from_id, to_id FROM relationships WHERE source_project_id = ?",
        (project_id,)
    ).fetchall()

    entity_ids = set()
    for r in rels:
        entity_ids.add(r["from_id"]); entity_ids.add(r["to_id"])

    materials, processes, perfs = [], [], []
    if entity_ids:
        ph = ",".join("?" * len(entity_ids))
        for e in db.execute(
            f"SELECT * FROM entities WHERE id IN ({ph})", list(entity_ids)
        ).fetchall():
            ed = dict(e)
            try:
                ed["properties"] = json.loads(ed.get("properties") or "{}")
            except Exception:
                ed["properties"] = {}
            if ed["type"] == "material":
                materials.append(ed)
            elif ed["type"] == "process":
                processes.append(ed)
            elif ed["type"] == "performance":
                perfs.append(ed)
    db.close()

    # Generate
    pdf_bytes = generate_pdf(project, materials, processes, perfs)
    pptx_bytes = generate_pptx(project, materials, processes, perfs)

    pdf_path.write_bytes(pdf_bytes)
    pptx_path.write_bytes(pptx_bytes)
    return pdf_path, pptx_path


def index_file_chunks_to_db():
    """For every project, generate the file content and insert ontology-grounded
    chunks into the chunks/chunks_fts tables so GraphRAG queries hit them.

    This runs after init_db. The actual files are generated lazily on download
    request — the chunk text is what the LLM/retrieval needs.
    """
    db = sqlite3.connect(_get_db_path())
    db.row_factory = sqlite3.Row

    # Already indexed? skip
    existing = db.execute(
        "SELECT COUNT(*) FROM chunks WHERE id LIKE '%-pdf-%' OR id LIKE '%-pptx-%'"
    ).fetchone()[0]
    if existing > 0:
        db.close()
        print(f"[files] {existing} synthetic file chunks already indexed")
        return existing

    chunks_to_insert = []
    project_ids = [r[0] for r in db.execute("SELECT id FROM projects").fetchall()]

    for pid in project_ids:
        prow = db.execute("SELECT * FROM projects WHERE id = ?", (pid,)).fetchone()
        project = dict(prow)
        for f in ("materials_used", "processes_used", "principal_investigators",
                  "industry_partners"):
            try:
                project[f] = json.loads(project.get(f) or "[]")
            except Exception:
                project[f] = []

        # Get entities
        rels = db.execute(
            "SELECT DISTINCT from_id, to_id FROM relationships WHERE source_project_id = ?",
            (pid,)
        ).fetchall()
        entity_ids = set()
        for r in rels:
            entity_ids.add(r["from_id"]); entity_ids.add(r["to_id"])

        materials, processes, perfs = [], [], []
        if entity_ids:
            ph = ",".join("?" * len(entity_ids))
            for e in db.execute(
                f"SELECT * FROM entities WHERE id IN ({ph})", list(entity_ids)
            ).fetchall():
                ed = dict(e)
                try:
                    ed["properties"] = json.loads(ed.get("properties") or "{}")
                except Exception:
                    ed["properties"] = {}
                if ed["type"] == "material":
                    materials.append(ed)
                elif ed["type"] == "process":
                    processes.append(ed)
                elif ed["type"] == "performance":
                    perfs.append(ed)

        # Build PDF section chunks
        for idx, mat in enumerate(materials[:4]):
            subclass = _classify_material(mat["name"])
            props_str = ", ".join(
                f"{k}={v}" for k, v in list(mat.get("properties", {}).items())[:4]
            )
            text = (
                f"FROM PDF FINAL REPORT (Section 2: Materials, ontology subclass: {subclass}). "
                f"{mat['name']} was selected for this project. Per the OWL-DL ontology, "
                f"this material is classified under {subclass}. "
                f"Vendor specifications and PERC/RURI characterization document the following "
                f"properties: {props_str}. "
                f"Materials selection follows NextFlex {project['project_call']} program criteria."
            )
            chunks_to_insert.append((
                f"{pid}-pdf-mat-{idx}-{mat['id'][:6]}", pid, "pdf:materials", 5,
                text, "public",
            ))

        for idx, proc in enumerate(processes[:3]):
            subclass = _classify_process(proc["name"])
            props_str = ", ".join(
                f"{k}={v}" for k, v in list(proc.get("properties", {}).items())[:5]
            )
            text = (
                f"FROM PDF FINAL REPORT (Section 3: Process & Manufacturing, "
                f"ontology subclass: {subclass}). The {proc['name']} process was operated with "
                f"the following parameters: {props_str}. "
                f"Per the OWL-DL ontology, this process is classified under {subclass}. "
                f"Process parameters were optimized through DOE on the {project['lead_institution']} "
                f"line for the target {project['focus_area']} application."
            )
            chunks_to_insert.append((
                f"{pid}-pdf-proc-{idx}-{proc['id'][:6]}", pid, "pdf:processes", 10,
                text, "public",
            ))

        for idx, perf in enumerate(perfs[:3]):
            subclass = _classify_performance(perf["name"])
            props_str = ", ".join(
                f"{k}={v}" for k, v in list(perf.get("properties", {}).items())[:5]
            )
            text = (
                f"FROM PDF FINAL REPORT (Section 4: Performance Results, "
                f"ontology subclass: {subclass}). Measured {perf['name']} for the integrated "
                f"demonstrator: {props_str}. "
                f"Per the OWL-DL ontology, this measurement is classified under {subclass}. "
                f"Statistical analysis across multiple independent samples confirms repeatability "
                f"within design tolerances."
            )
            chunks_to_insert.append((
                f"{pid}-pdf-perf-{idx}-{perf['id'][:6]}", pid, "pdf:results", 15,
                text, "public",
            ))

        # PPTX briefing chunks (executive summary + transition)
        partners = ", ".join(project.get("industry_partners", [])) or "TBD"
        pptx_summary = (
            f"FROM PPTX PROJECT BRIEFING (Slide 2: Executive Summary). "
            f"Project: {project['title']}. Lead: {project['lead_institution']}. "
            f"PIs: {', '.join(project.get('principal_investigators', []))}. "
            f"Funding: ${project['funding_amount']:,}. "
            f"TRL advancement: {project['trl_start']} to {project['trl_end']}. "
            f"Abstract: {project['abstract'][:240]}."
        )
        chunks_to_insert.append((
            f"{pid}-pptx-summary", pid, "pptx:executive_summary", 2,
            pptx_summary, "public",
        ))

        pptx_transition = (
            f"FROM PPTX PROJECT BRIEFING (Slide 7: Transition Pathway). "
            f"Industry partners: {partners}. "
            f"PEO alignment: {project.get('peo') or 'TBD'}. "
            f"Classification: {project['classification']}. "
            f"Results entered into NextFlex Secure AI/ML Knowledge Base for cross-project "
            f"reasoning. Queryable via GraphRAG interface."
        )
        chunks_to_insert.append((
            f"{pid}-pptx-transition", pid, "pptx:transition", 7,
            pptx_transition, "public",
        ))

    # Bulk insert
    db.executemany(
        "INSERT INTO chunks (id, project_id, section, page, text, classification) "
        "VALUES (?, ?, ?, ?, ?, ?)",
        chunks_to_insert,
    )
    # Insert into FTS
    db.executemany(
        "INSERT INTO chunks_fts (id, text, project_id, section) VALUES (?, ?, ?, ?)",
        [(c[0], c[4], c[1], c[2]) for c in chunks_to_insert],
    )
    db.commit()
    n = len(chunks_to_insert)
    db.close()
    print(f"[files] Indexed {n} synthetic file chunks into FTS")
    return n


if __name__ == "__main__":
    n = index_file_chunks_to_db()
    print(f"Total file chunks indexed: {n}")
    # Quick test: generate one file
    db = sqlite3.connect(_get_db_path())
    pid = db.execute("SELECT id FROM projects LIMIT 1").fetchone()[0]
    db.close()
    pdf_path, pptx_path = get_or_generate_files(pid)
    print(f"Test PDF: {pdf_path} ({pdf_path.stat().st_size:,} bytes)")
    print(f"Test PPTX: {pptx_path} ({pptx_path.stat().st_size:,} bytes)")
