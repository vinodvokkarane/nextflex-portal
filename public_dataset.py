"""
Public NextFlex Dataset indexer.

Loads the real NextFlex acknowledged papers and funded assets into the
existing portal tables so they're queryable through the same UI:
  - GraphRAG retrieves from real PDF/PPTX content
  - Drill-downs work via project-style records
  - Files are downloadable through /api/public-dataset/files/{id}

Schema additions to nextflex.db:
  - public_assets table: holds metadata for each real PDF/PPTX file
  - public_assets_fts: FTS5 index on title + abstract + extracted text
  - chunks: extended with classification='public_dataset' for the
    actual paper text content (so existing GraphRAG endpoint hits them)

All records carry classification='public' so every persona sees them.
"""

import csv
import json
import re
import sqlite3
from pathlib import Path

DB_PATH = Path(__file__).parent / "nextflex.db"
DATA_ROOT = Path(__file__).parent / "public_data"
PAPERS_DIR = DATA_ROOT / "nextflex_acknowledged_papers_2016_2026"
ASSETS_DIR = DATA_ROOT / "nextflex_assets"


PUBLIC_SCHEMA = """
CREATE TABLE IF NOT EXISTS public_assets (
    id TEXT PRIMARY KEY,
    category TEXT NOT NULL,        -- paper | patent | webinar | project_report | pptx | podcast | presentation
    subtype TEXT,                  -- patent_pdf | webinar_pdf | etc.
    title TEXT NOT NULL,
    year INTEGER,
    authors TEXT,                  -- JSON array
    project_call TEXT,             -- e.g. PC 1.0 (best-effort inferred)
    funding_acknowledgment TEXT,
    agreement_numbers TEXT,
    file_path TEXT,                -- relative path under public_data/
    file_size_bytes INTEGER,
    pages INTEGER,
    source_url TEXT,
    public_access TEXT,
    download_status TEXT,
    has_local_file INTEGER DEFAULT 0,
    abstract TEXT,                 -- first ~500 chars of extracted text
    full_text_chars INTEGER,
    created_at TEXT DEFAULT CURRENT_TIMESTAMP
);

CREATE INDEX IF NOT EXISTS idx_pub_cat ON public_assets (category);
CREATE INDEX IF NOT EXISTS idx_pub_year ON public_assets (year);
CREATE INDEX IF NOT EXISTS idx_pub_pc ON public_assets (project_call);
CREATE INDEX IF NOT EXISTS idx_pub_local ON public_assets (has_local_file);

CREATE VIRTUAL TABLE IF NOT EXISTS public_assets_fts USING fts5(
    id UNINDEXED, title, authors, abstract, agreement_numbers, project_call
);
"""


def _safe_id(category: str, filename: str) -> str:
    """Stable, FS-safe ID for an asset."""
    base = Path(filename).stem.lower()
    base = re.sub(r"[^a-z0-9]+", "_", base).strip("_")[:80]
    return f"pub-{category}-{base}"


def _infer_project_call(text: str, default: str = "") -> str:
    """Try to guess a Project Call number from the text."""
    if not text:
        return default
    # Look for explicit "PC X.X" or "Project Call X.X" mentions
    m = re.search(r"\b(PC|Project Call|OPC)[\s\-]*?(\d{1,2}(?:\.\d)?)\b", text)
    if m:
        n = m.group(2)
        return f"PC {n}"
    return default


def _extract_pdf_text(path: Path, max_pages: int = 40) -> tuple[str, int]:
    """Extract text from a PDF. Returns (text, pages_total)."""
    try:
        import pypdf
    except ImportError:
        return "", 0
    try:
        reader = pypdf.PdfReader(str(path))
        pages_total = len(reader.pages)
        chunks = []
        for i, page in enumerate(reader.pages[:max_pages]):
            try:
                chunks.append(page.extract_text() or "")
            except Exception:
                pass
        return "\n".join(chunks), pages_total
    except Exception as e:
        print(f"[public-data] pdf extract failed for {path.name}: {e}", flush=True)
        return "", 0


def _extract_pptx_text(path: Path) -> tuple[str, int]:
    """Extract text from a PPTX."""
    try:
        from pptx import Presentation
    except ImportError:
        return "", 0
    try:
        prs = Presentation(str(path))
        chunks = []
        n_slides = 0
        for slide in prs.slides:
            n_slides += 1
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            slide_text.append(t)
            if slide_text:
                chunks.append(f"[Slide {n_slides}]\n" + "\n".join(slide_text))
        return "\n\n".join(chunks), n_slides
    except Exception as e:
        print(f"[public-data] pptx extract failed for {path.name}: {e}", flush=True)
        return "", 0


def _chunk_text(text: str, chunk_chars: int = 1500, overlap: int = 150):
    """Split text into chunks suitable for retrieval."""
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return []
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_chars, len(text))
        # Prefer a sentence boundary near the end
        if end < len(text):
            for p in (".", "!", "?", "\n"):
                idx = text.rfind(p, start + chunk_chars // 2, end)
                if idx > 0:
                    end = idx + 1
                    break
        chunks.append(text[start:end].strip())
        start = end - overlap if end - overlap > start else end
    return [c for c in chunks if len(c) > 50]


def index_public_dataset():
    """Idempotent: index the real NextFlex papers + funded assets."""
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    conn.executescript(PUBLIC_SCHEMA)

    # Skip if already indexed
    existing = conn.execute(
        "SELECT COUNT(*) FROM public_assets"
    ).fetchone()[0]
    if existing > 0:
        print(f"[public-data] {existing} assets already indexed, skipping", flush=True)
        conn.close()
        return existing

    chunk_rows = []  # (id, project_id, section, page, text, classification)

    # ─── Papers ────────────────────────────────────────────────
    papers_manifest = PAPERS_DIR / "manifest.csv"
    if papers_manifest.exists():
        n_papers = 0
        with papers_manifest.open(newline="", encoding="utf-8") as f:
            for row in csv.DictReader(f):
                year = int(row.get("year") or 0) or None
                title = (row.get("title") or "").strip()
                if not title:
                    continue
                folder = (row.get("folder") or "").strip()
                filename = (row.get("filename") or "").strip()
                rel_path = f"nextflex_acknowledged_papers_2016_2026/{folder}/{filename}" if filename else ""
                local_path = PAPERS_DIR / folder / filename if filename else None
                size = int(row.get("size_bytes") or 0)
                agreements = (row.get("agreements") or "").strip()
                ack_excerpt = (row.get("ack_or_funding_excerpt") or "").strip()

                # Extract PDF text if file present
                full_text = ""
                pages = 0
                if local_path and local_path.exists():
                    full_text, pages = _extract_pdf_text(local_path)

                abstract = (full_text[:600] if full_text else ack_excerpt[:600]).strip()

                aid = _safe_id("paper", filename or title)
                inferred_pc = _infer_project_call(full_text or ack_excerpt)

                conn.execute(
                    """INSERT OR REPLACE INTO public_assets
                       (id, category, subtype, title, year, project_call,
                        funding_acknowledgment, agreement_numbers,
                        file_path, file_size_bytes, pages,
                        public_access, download_status, has_local_file,
                        abstract, full_text_chars)
                       VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)""",
                    (aid, "paper", "academic_paper", title, year, inferred_pc,
                     ack_excerpt, agreements, rel_path, size, pages,
                     "public", row.get("status", "included_public_pdf"),
                     1 if local_path and local_path.exists() else 0,
                     abstract, len(full_text)),
                )
                n_papers += 1

                # Chunk paper text for GraphRAG
                if full_text:
                    for ci, ct in enumerate(_chunk_text(full_text)):
                        chunk_rows.append((
                            f"{aid}-c{ci:03d}", aid, "public:paper",
                            ci + 1, ct, "public_dataset",
                        ))
        print(f"[public-data] {n_papers} papers indexed", flush=True)

    # ─── Funded assets ─────────────────────────────────────────
    assets_manifest = DATA_ROOT / "nextflex_assets_included_public_files.csv"
    if assets_manifest.exists():
        n_assets = 0
        with assets_manifest.open(newline="", encoding="utf-8") as f:
            for row in csv.DictReader(f):
                category = (row.get("category") or "").strip()
                subtype = (row.get("subtype") or "").strip()
                title = (row.get("title") or "").strip()
                year_raw = (row.get("year") or "").strip()
                try:
                    year = int(year_raw) if year_raw else None
                except ValueError:
                    year = None
                rel_path = (row.get("downloaded_filename") or "").strip()
                file_exists = (row.get("file_exists") or "").strip().lower() == "yes"
                size = int(row.get("file_size_bytes") or 0)
                source_url = (row.get("source_url") or "").strip()
                agreements = (row.get("agreement_numbers") or "").strip()
                evidence = (row.get("evidence") or "").strip()
                public_access = (row.get("public_access") or "").strip()
                download_status = (row.get("download_status") or "").strip()

                if not title:
                    continue

                local_path = DATA_ROOT / rel_path if rel_path else None
                full_text = ""
                pages = 0
                if file_exists and local_path and local_path.exists():
                    if rel_path.lower().endswith(".pdf"):
                        full_text, pages = _extract_pdf_text(local_path)
                    elif rel_path.lower().endswith(".pptx"):
                        full_text, pages = _extract_pptx_text(local_path)

                abstract = (full_text[:600] if full_text else evidence[:600]).strip()

                aid = _safe_id(category or "asset", Path(rel_path).name or title)
                inferred_pc = _infer_project_call(full_text)

                conn.execute(
                    """INSERT OR REPLACE INTO public_assets
                       (id, category, subtype, title, year, project_call,
                        funding_acknowledgment, agreement_numbers,
                        file_path, file_size_bytes, pages, source_url,
                        public_access, download_status, has_local_file,
                        abstract, full_text_chars)
                       VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)""",
                    (aid, category, subtype, title, year, inferred_pc,
                     evidence, agreements, rel_path, size, pages, source_url,
                     public_access, download_status,
                     1 if file_exists else 0,
                     abstract, len(full_text)),
                )
                n_assets += 1

                # Chunk asset text
                if full_text:
                    for ci, ct in enumerate(_chunk_text(full_text)):
                        chunk_rows.append((
                            f"{aid}-c{ci:03d}", aid, f"public:{category}",
                            ci + 1, ct, "public_dataset",
                        ))
        print(f"[public-data] {n_assets} funded assets indexed", flush=True)

    # ─── Insert all chunks into the same chunks table used by GraphRAG ─
    if chunk_rows:
        # Use INSERT OR IGNORE in case re-indexing
        conn.executemany(
            "INSERT OR IGNORE INTO chunks "
            "(id, project_id, section, page, text, classification) "
            "VALUES (?, ?, ?, ?, ?, ?)",
            chunk_rows,
        )
        # Mirror into chunks_fts
        conn.executemany(
            "INSERT OR IGNORE INTO chunks_fts (id, text, project_id, section) "
            "VALUES (?, ?, ?, ?)",
            [(r[0], r[4], r[1], r[2]) for r in chunk_rows],
        )
        print(f"[public-data] {len(chunk_rows)} text chunks indexed for GraphRAG",
              flush=True)

    # ─── Build assets FTS ──────────────────────────────────────
    rows = conn.execute(
        "SELECT id, title, COALESCE(authors,''), COALESCE(abstract,''), "
        "       COALESCE(agreement_numbers,''), COALESCE(project_call,'') "
        "FROM public_assets"
    ).fetchall()
    conn.executemany(
        "INSERT INTO public_assets_fts "
        "(id, title, authors, abstract, agreement_numbers, project_call) "
        "VALUES (?, ?, ?, ?, ?, ?)",
        [tuple(r) for r in rows],
    )

    conn.commit()
    total = conn.execute("SELECT COUNT(*) FROM public_assets").fetchone()[0]
    chunks = conn.execute(
        "SELECT COUNT(*) FROM chunks WHERE classification='public_dataset'"
    ).fetchone()[0]
    conn.close()
    print(f"[public-data] DONE: {total} assets, {chunks} chunks indexed",
          flush=True)
    return total


if __name__ == "__main__":
    n = index_public_dataset()
    print(f"Indexed {n} public assets")
